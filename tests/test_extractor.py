from datetime import date
from pathlib import Path
from pptx import Presentation

from pptx_schedule import ExtractionConfig, extract_project_info


def _create_sample_pptx(path: Path) -> None:
    presentation = Presentation()
    blank_slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_slide_layout)

    textbox = slide.shapes.add_textbox(left=100000, top=100000, width=4000000, height=500000)
    textbox.text = "專案名稱: 測試專案"

    rows, cols = 3, 2
    table_shape = slide.shapes.add_table(rows, cols, left=100000, top=1500000, width=5000000, height=2000000)
    table = table_shape.table
    table.cell(0, 0).text = "Schedule"
    table.cell(0, 1).text = "Item"
    table.cell(1, 0).text = "2024/05/01"
    table.cell(1, 1).text = "Kickoff"
    table.cell(2, 0).text = "2024年05月10日"
    table.cell(2, 1).text = "Review"

    presentation.save(path)


def test_extract_project_info(tmp_path: Path) -> None:
    pptx_path = tmp_path / "demo.pptx"
    _create_sample_pptx(pptx_path)

    result = extract_project_info(pptx_path)

    assert result.project_name == "測試專案"
    dates = sorted(entry.value for entry in result.schedule_entries)
    assert dates == [date(2024, 5, 1), date(2024, 5, 10)]


def test_fallback_to_filename(tmp_path: Path) -> None:
    pptx_path = tmp_path / "no_text.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(pptx_path)

    result = extract_project_info(pptx_path, config=ExtractionConfig(project_name_patterns=[]))
    assert result.project_name == "no_text"
    assert result.schedule_entries == []
