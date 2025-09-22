"""Functions that parse PPTX presentations and extract scheduling information."""

from __future__ import annotations

import re
from datetime import date
from pathlib import Path
from typing import Iterator, List, Optional, Sequence, Tuple

from dateutil import parser
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .models import ExtractionConfig, ExtractedProject, ScheduleEntry

DATE_TOKEN_PATTERNS: Sequence[re.Pattern[str]] = (
    re.compile(r"\b\d{4}[/-]\d{1,2}[/-]\d{1,2}\b"),
    re.compile(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b"),
    re.compile(r"\b\d{1,2}[/-]\d{1,2}\b"),
)
CHINESE_DATE_PATTERN = re.compile(
    r"(?P<year>\d{4})年(?P<month>\d{1,2})月(?P<day>\d{1,2})日"
)


def extract_project_info(
    pptx_path: Path | str, config: Optional[ExtractionConfig] = None
) -> ExtractedProject:
    """Extract the project name and schedule entries from the provided PPTX file."""

    cfg = config or ExtractionConfig()
    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX file not found: {path}")

    presentation = Presentation(str(path))
    texts: List[str] = []
    schedule_entries: List[ScheduleEntry] = []
    seen_entries: set[Tuple[int, date, str]] = set()

    for slide_index, slide in enumerate(presentation.slides):
        for text in _iter_slide_text(slide.shapes):
            cleaned = _clean_text(text)
            if cleaned:
                texts.append(cleaned)

        for table in _iter_slide_tables(slide.shapes):
            if not _table_matches_keywords(table, cfg.schedule_table_keywords):
                continue

            for entry in _extract_schedule_entries_from_table(
                table, slide_index=slide_index, dayfirst=cfg.dayfirst
            ):
                entry_key = (entry.slide_index, entry.value, entry.cell_text)
                if entry_key not in seen_entries:
                    seen_entries.add(entry_key)
                    schedule_entries.append(entry)

    project_name = _resolve_project_name(texts, cfg, fallback=path.stem)

    return ExtractedProject(
        project_name=project_name,
        schedule_entries=schedule_entries,
        source_path=path,
    )


def _iter_slide_text(shapes) -> Iterator[str]:
    for shape in shapes:
        yield from _iter_shape_text(shape)


def _iter_shape_text(shape) -> Iterator[str]:
    if getattr(shape, "has_text_frame", False):
        yield shape.text
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:  # type: ignore[attr-defined]
            yield from _iter_shape_text(sub_shape)


def _iter_slide_tables(shapes) -> Iterator:
    for shape in shapes:
        yield from _iter_shape_tables(shape)


def _iter_shape_tables(shape) -> Iterator:
    if getattr(shape, "has_table", False):
        yield shape.table
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:  # type: ignore[attr-defined]
            yield from _iter_shape_tables(sub_shape)


def _table_matches_keywords(table, keywords: Sequence[str]) -> bool:
    if not keywords:
        return True

    joined = " ".join(
        _clean_text(cell.text)
        for row in table.rows
        for cell in row.cells
        if cell.text
    )
    haystack = joined.lower()
    return any(keyword.lower() in haystack for keyword in keywords)


def _extract_schedule_entries_from_table(table, *, slide_index: int, dayfirst: bool) -> Iterator[ScheduleEntry]:
    for row in table.rows:
        for cell in row.cells:
            text = _clean_text(cell.text)
            if not text:
                continue
            for value in _extract_dates_from_text(text, dayfirst=dayfirst):
                yield ScheduleEntry(slide_index=slide_index, cell_text=text, value=value)


def _extract_dates_from_text(text: str, *, dayfirst: bool) -> Iterator[date]:
    seen: set[Tuple[int, int, int]] = set()
    for pattern in DATE_TOKEN_PATTERNS:
        for match in pattern.finditer(text):
            candidate = match.group(0)
            parsed = _parse_date(candidate, dayfirst=dayfirst)
            if parsed:
                key = (parsed.year, parsed.month, parsed.day)
                if key not in seen:
                    seen.add(key)
                    yield parsed

    for match in CHINESE_DATE_PATTERN.finditer(text):
        year = int(match.group("year"))
        month = int(match.group("month"))
        day = int(match.group("day"))
        key = (year, month, day)
        if key not in seen:
            seen.add(key)
            yield date(year, month, day)


def _parse_date(candidate: str, *, dayfirst: bool) -> Optional[date]:
    try:
        parsed = parser.parse(candidate, fuzzy=False, dayfirst=dayfirst)
    except (ValueError, OverflowError):
        return None
    return parsed.date()


def _resolve_project_name(
    texts: Sequence[str], config: ExtractionConfig, *, fallback: str
) -> str:
    patterns = [re.compile(pattern, re.IGNORECASE) for pattern in config.project_name_patterns]

    for text in texts:
        for pattern in patterns:
            match = pattern.search(text)
            if match:
                value = match.groupdict().get("value")
                if value is None and match.groups():
                    value = match.group(match.lastindex or 0)
                if value:
                    cleaned = _normalize_name(value)
                    if cleaned:
                        return cleaned

    for text in texts:
        lowered = text.lower()
        for keyword in config.project_name_keywords:
            keyword_lower = keyword.lower()
            if keyword_lower in lowered:
                extracted = _strip_keyword(text, keyword_lower)
                if extracted:
                    return extracted

    for text in texts:
        cleaned = _normalize_name(text)
        if cleaned:
            return cleaned

    return _normalize_name(fallback) or fallback


def _strip_keyword(text: str, keyword_lower: str) -> Optional[str]:
    lowered = text.lower()
    index = lowered.find(keyword_lower)
    if index < 0:
        return None
    result = text[index + len(keyword_lower) :]
    result = result.lstrip(" :：-\n\t")
    return _normalize_name(result)


def _normalize_name(value: str) -> str:
    cleaned = _clean_text(value)
    return cleaned


def _clean_text(value: str | None) -> str:
    if not value:
        return ""
    collapsed = re.sub(r"\s+", " ", value)
    return collapsed.strip(" \n\t:\uff1a")
